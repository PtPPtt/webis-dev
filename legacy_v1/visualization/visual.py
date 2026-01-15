import argparse
import os
import pathlib
import re
import uuid
import webbrowser
import json
from typing import List, Dict, Any
import requests
from dotenv import load_dotenv

# PPT/PDF Libraries
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as ReportLabImage
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
except ImportError:
    print("Warning: python-pptx or reportlab not installed. PDF/PPT export will not work.")

# 加载环境变量
load_dotenv(".env.local")

# 获取 SILICONFLOW_API_KEY
API_KEY = os.getenv("SILICONFLOW_API_KEY")
if not API_KEY:
    raise RuntimeError("缺少 SILICONFLOW_API_KEY，请在 .env.local 文件中设置。")


def _strip_markdown_fences(text: str) -> str:
    stripped = text.strip()
    m = re.search(r"```(?:html)?\s*(.*?)\s*```", stripped, flags=re.DOTALL | re.IGNORECASE)
    if m:
        return m.group(1).strip()
    if stripped.startswith("```"):
        lines = stripped.splitlines()[1:]
        stripped = "\n".join(lines)
        stripped = re.sub(r"\n```\s*$", "", stripped)
        return stripped.strip()
    return stripped


def _sanitize_html(html: str) -> str:
    s = _strip_markdown_fences(html)

    s = s.replace("“", '"').replace("”", '"').replace("‘", "'").replace("’", "'")

    # Fix common model glitch patterns like '<< section ...' and '< / div >'
    s = re.sub(r"<<\s*/\s*([A-Za-z])", r"</\1", s)
    s = re.sub(r"<<\s*([A-Za-z])", r"<\1", s)

    s = re.sub(r"<\s*/\s*([A-Za-z][\w:-]*)\s*>", r"</\1>", s)
    s = re.sub(r"<\s+([!/A-Za-z])", r"<\1", s)

    attrs = r"(?:class|id|href|src|style|lang|type|rel|name|content|crossorigin|role)"
    s = re.sub(rf"(<\s*/?\s*)([A-Za-z][\w:-]*?)({attrs})\s*=", r"\1\2 \3=", s)
    s = re.sub(r"(<\s*/?\s*)([A-Za-z][\w:-]*?)(aria-[\w-]+)\s*=", r"\1\2 \3=", s)

    s = re.sub(r'(\s[\w:-]+)\s*=\s*"', r'\1="', s)

    def _normalize_style_block(m: re.Match) -> str:
        style = m.group(1)
        style = (
            style.replace("：", ":")
            .replace("；", ";")
            .replace("（", "(")
            .replace("）", ")")
            .replace("，", ",")
        )
        style = re.sub(r"\b(var|calc|rgba|rgb)\s*\(\s*", r"\1(", style)
        style = re.sub(r"\(\s*--", "(--", style)
        return f"<style>{style}</style>"

    s = re.sub(r"<style[^>]*>(.*?)</style>", _normalize_style_block, s, flags=re.DOTALL | re.IGNORECASE)
    return s.strip()


def _find_html_issues(html: str) -> List[str]:
    issues: List[str] = []
    if "<!DOCTYPE html" not in html and "<!doctype html" not in html:
        issues.append("missing_doctype")
    if "</html>" not in html.lower():
        issues.append("missing_html_close")
    if re.search(r"<<\s*/?\s*[A-Za-z]", html):
        issues.append("double_angle_brackets_in_tags")
    if re.search(r"<\s+/?\s*[A-Za-z]", html):
        issues.append("whitespace_after_lt")
    if re.search(r"</\s+[A-Za-z]", html):
        issues.append("whitespace_in_close_tag")
    if "```" in html:
        issues.append("markdown_fence_present")
    if any(ch in html for ch in ["“", "”", "‘", "’"]):
        issues.append("curly_quotes_present")
    if any(ch in html for ch in ["：", "；", "（", "）"]):
        issues.append("fullwidth_punct_present")
    return issues


def _repair_html_with_model(url: str, headers: dict, broken_html: str) -> str:
    repair_payload = {
        "model": "Qwen/Qwen2.5-Coder-32B-Instruct",
        "messages": [
            {
                "role": "system",
                "content": (
                    "你是一个专业的网页开发者。你的任务是修复给定的HTML，使其成为完整、格式正确、可渲染的HTML5文档。\n"
                    "要求：\n"
                    "1) 以<!DOCTYPE html>开头，以</html>结束\n"
                    "2) 修复类似 '<<section'、'< / div >' 这类标签错误\n"
                    "3) 修复中文引号“”与全角标点（：；（））导致的HTML/CSS语法错误，统一为英文半角\n"
                    "4) 不要输出```html这类Markdown代码块包裹\n"
                    "5) 只输出修复后的HTML代码，不要额外说明"
                ),
            },
            {"role": "user", "content": broken_html},
        ],
        "stream": False,
        "max_tokens": 10000,
        "temperature": 0.6,
        "top_p": 0.7,
        "top_k": 50,
        "n": 1,
        "response_format": {"type": "text"},
    }

    repair_resp = requests.post(url, json=repair_payload, headers=headers, timeout=120)
    repair_resp.raise_for_status()
    repair_json = repair_resp.json()
    return repair_json["choices"][0]["message"]["content"]


def _generate_slide_content(user_content: str) -> List[Dict[str, Any]]:
    """
    Uses LLM to generate structured content for slides (Title, Bullet points).
    """
    url = "https://api.siliconflow.cn/v1/chat/completions"
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}

    system_prompt = """You are a presentation expert. 
    Convert the given content into a structured JSON format for a presentation.
    Output a JSON object with a key "slides" which is a list of objects.
    Each slide object must have:
    - "title": string
    - "content": list of strings (bullet points)
    
    Example:
    {
      "slides": [
        {"title": "Introduction", "content": ["Point 1", "Point 2"]},
        {"title": "Details", "content": ["Detail A", "Detail B"]}
      ]
    }
    Only output valid JSON."""

    payload = {
        "model": "deepseek-ai/DeepSeek-V3", # Use a fast model
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": f"Content to convert:\n{user_content}"},
        ],
        "response_format": {"type": "json_object"},
        "max_tokens": 4000,
    }

    try:
        response = requests.post(url, json=payload, headers=headers, timeout=60)
        response.raise_for_status()
        data = response.json()
        content = data["choices"][0]["message"]["content"]
        return json.loads(content).get("slides", [])
    except Exception as e:
        print(f"Error generating slide content: {e}")
        return []


def _create_ppt(slides: List[Dict[str, Any]], output_path: str):
    """Generates a PowerPoint file from slide data."""
    try:
        prs = Presentation()
        for slide_data in slides:
            slide_layout = prs.slide_layouts[1] # Bullet slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title = slide.shapes.title
            title.text = slide_data.get("title", "Untitled")
            
            # Content
            tf = slide.placeholders[1].text_frame
            for point in slide_data.get("content", []):
                p = tf.add_paragraph()
                p.text = point
                p.level = 0

        prs.save(output_path)
        print(f"PPT generated: {output_path}")
    except Exception as e:
        print(f"Failed to create PPT: {e}")


def _create_pdf(slides: List[Dict[str, Any]], output_path: str):
    """Generates a PDF file from slide data."""
    try:
        # Register Chinese font support
        try:
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.cidfonts import UnicodeCIDFont
            pdfmetrics.registerFont(UnicodeCIDFont('STSong-Light'))
            font_name = 'STSong-Light'
        except Exception:
            font_name = 'Helvetica'
            print("Warning: Chinese font not found, PDF may not render Chinese characters correctly.")

        doc = SimpleDocTemplate(output_path, pagesize=letter)
        styles = getSampleStyleSheet()
        
        # Update styles to use the font
        title_style = styles['Title']
        title_style.fontName = font_name
        body_style = styles['BodyText']
        body_style.fontName = font_name
        
        story = []

        for slide in slides:
            # Title
            story.append(Paragraph(slide.get("title", "Untitled"), title_style))
            story.append(Spacer(1, 12))
            
            # Content
            for point in slide.get("content", []):
                # Simple bullet point simulation
                text = f"• {point}"
                story.append(Paragraph(text, body_style))
                story.append(Spacer(1, 6))
            
            story.append(Spacer(1, 24)) # Space between slides/sections

        doc.build(story)
        print(f"PDF generated: {output_path}")
    except Exception as e:
        print(f"Failed to create PDF: {e}")


def _resolve_path(input_path: str) -> str:
    if not input_path:
        raise ValueError("Missing input file path.")
    path = pathlib.Path(input_path)
    if path.is_absolute():
        return str(path)
    project_root = pathlib.Path(__file__).resolve().parents[1]
    return str(project_root / input_path.lstrip("\\/"))


def main(argv: List[str] | None = None) -> None:
    parser = argparse.ArgumentParser(description="Generate HTML visualization from input text.")
    parser.add_argument("input_path", nargs="?", help="Input file path.")
    parser.add_argument("--input", dest="input_flag", help="Input file path.")
    parser.add_argument("--pdf", action="store_true", help="Generate PDF output")
    parser.add_argument("--ppt", action="store_true", help="Generate PPT output")
    args = parser.parse_args(argv)

    input_path = (args.input_flag or args.input_path or "").strip()
    if not input_path:
        input_path = input("请输入输入文件路径: ").strip()

    input_file_path = _resolve_path(input_path)

    with open(input_file_path, "r", encoding="utf-8") as file:
        user_content = file.read()

    # 1. Generate HTML (Existing Logic)
    print("Generating HTML visualization...")
    url = "https://api.siliconflow.cn/v1/chat/completions"
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}

    payload = {
        "model": "Pro/deepseek-ai/DeepSeek-R1",
        "messages": [
            {
                "role": "system",
                "content": """你是一个专业的网页开发者。请生成完整、格式正确、可渲染的HTML5代码。

重要要求：
1) 必须生成完整HTML文档：以<!DOCTYPE html>开头，以</html>结束
2) 所有标签必须正确闭合；不要输出 <p >、</ div > 这类“尖括号附近有空格”的标签
3) 标签名/属性名必须为英文半角；属性值必须使用英文双引号 " 包裹（不要用中文引号“”）
4) CSS/JS 必须使用英文半角标点（: ; ( ) , .），不要出现中文全角标点（例如：：；（））
5) 不要输出类似 <<div>、<<footer> 这种重复尖括号的标签
6) 不要输出 ```html 这类Markdown代码块包裹
7) 只输出HTML代码，不要额外说明；输出前请自检并修复以上问题""",
            },
            {
                "role": "user",
                "content": f"""请为我创建美观的HTML网页，要求：

设计规范：
1. 完整HTML5结构
2. 现代化浅色UI配色
3. 合理的内容分区，网页结构不要过于简单
4. 优雅的排版
5. 内容来源只能使用“基于Webis生成”标志
6. 可以添加合理配图

内容来源（合理总结使用）：
{user_content}

重要：只输出HTML代码，不要额外说明""",
            },
        ],
        "stream": False,
        "max_tokens": 20000,
        "min_p": 0.05,
        "stop": None,
        "temperature": 0.3,
        "top_p": 0.9,
        "top_k": 50,
        "frequency_penalty": 0.0,
        "n": 1,
        "response_format": {"type": "text"},
    }

    response = requests.post(url, json=payload, headers=headers, timeout=300)
    try:
        response.raise_for_status()
    except requests.HTTPError:
        print("HTTP error from API:", response.text)
        raise

    response_json = response.json()
    if "error" in response_json:
        print("API error:", response_json["error"])
        raise KeyError("choices")
    if "choices" not in response_json:
        print("Unexpected API response:", response_json)
        raise KeyError("choices")

    raw_html_content = response_json["choices"][0]["message"]["content"]
    html_content = _sanitize_html(raw_html_content)

    issues = _find_html_issues(html_content)
    if issues and os.environ.get("WEBIS_REPAIR_HTML") in {"1", "true", "TRUE", "yes", "YES"}:
        repaired = _repair_html_with_model(url, headers, html_content)
        html_content = _sanitize_html(repaired)
        issues = _find_html_issues(html_content)

    if issues:
        print("Generated HTML looks malformed:", issues)

    random_id = uuid.uuid4().hex
    output_folder = os.path.join("pipeline_outputs", "output_HTML")
    os.makedirs(output_folder, exist_ok=True)

    html_filename = f"generated_page_{random_id}.html"
    html_path = os.path.join(output_folder, html_filename)

    with open(html_path, "w", encoding="utf-8") as file:
        file.write(html_content)
    webbrowser.open(pathlib.Path(html_path).resolve().as_uri())

    print(f"HTML文件已成功生成！文件名: {html_path}")

    # 2. Generate PDF/PPT if requested
    if args.pdf or args.ppt:
        print("Generating structured content for PDF/PPT...")
        slides_data = _generate_slide_content(user_content)
        
        if args.ppt:
            ppt_filename = f"generated_presentation_{random_id}.pptx"
            ppt_path = os.path.join(output_folder, ppt_filename)
            _create_ppt(slides_data, ppt_path)
            
        if args.pdf:
            pdf_filename = f"generated_document_{random_id}.pdf"
            pdf_path = os.path.join(output_folder, pdf_filename)
            _create_pdf(slides_data, pdf_path)


if __name__ == "__main__":
    main()